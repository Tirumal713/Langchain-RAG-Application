import os
from typing import List, Dict
import PyPDF2
import pptx
import openpyxl
from langchain.text_splitter import RecursiveCharacterTextSplitter

class DocumentProcessor:
    @staticmethod
    def read_pdf(file_path: str) -> List[Dict[str, str]]:
        chunks = []
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page_num, page in enumerate(pdf_reader.pages, 1):
                # Extract text for the entire page
                text = page.extract_text()
                
                # Create a chunk for each page
                chunks.append({
                    'text': text,
                    'source': f'{os.path.basename(file_path)}_page_{page_num}',
                    'page': page_num
                })
        return chunks

    @staticmethod
    def read_pptx(file_path: str) -> List[Dict[str, str]]:
        chunks = []
        prs = pptx.Presentation(file_path)
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            
            # Join all text from the slide
            full_text = ' '.join(slide_text)
            
            # Create a chunk for each slide
            chunks.append({
                'text': full_text,
                'source': f'{os.path.basename(file_path)}_slide_{slide_num}',
                'page': slide_num
            })
        return chunks

    @staticmethod
    def read_excel(file_path: str) -> List[Dict[str, str]]:
        chunks = []
        wb = openpyxl.load_workbook(file_path)
        
        # Collect text from all sheets
        all_sheet_data = []
        for sheet_num, sheet in enumerate(wb.worksheets, 1):
            sheet_data = []
            for row in sheet.iter_rows(values_only=True):
                row_text = ' '.join(str(cell) for cell in row if cell is not None)
                sheet_data.append(row_text)
            
            # Join sheet data
            sheet_full_text = ' '.join(sheet_data)
            all_sheet_data.append(sheet_full_text)
        
        # Create a single chunk for the entire workbook
        full_workbook_text = ' '.join(all_sheet_data)
        chunks.append({
            'text': full_workbook_text,
            'source': f'{os.path.basename(file_path)}_workbook',
            'page': 1
        })
        
        return chunks


